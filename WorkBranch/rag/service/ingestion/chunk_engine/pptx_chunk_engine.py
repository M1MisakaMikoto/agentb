from __future__ import annotations

from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base_chunk_engine import BaseChunkEngine
from .file_profile import FileProfile


class PptxChunkEngine(BaseChunkEngine):
    """Extract text from modern .pptx (OOXML zip) via python-pptx, per slide."""

    name = "pptx"

    _PPTX_SUFFIX = {".pptx"}
    _PPTX_MIME_KEYWORDS = ("presentationml.presentation",)

    def can_handle(self, profile: FileProfile) -> float:
        if profile.extension in self._PPTX_SUFFIX:
            return 1.0
        lowered_mime = (profile.mime or "").lower()
        if any(keyword in lowered_mime for keyword in self._PPTX_MIME_KEYWORDS):
            return 0.95
        return 0.0

    def extract_text(self, file_path: Path) -> str:
        prs = Presentation(str(file_path))
        blocks: List[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            self._collect_shape_texts(slide.shapes, texts)
            body = "\n".join(t for t in texts if t.strip())
            if body:
                blocks.append(f"[slide:{idx}]\n{body}")
        return "\n\n".join(blocks)

    def _collect_shape_texts(self, shapes, out: List[str]) -> None:
        for shape in shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
                out.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    row_text = " | ".join([c for c in cells if c])
                    if row_text:
                        out.append(row_text)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._collect_shape_texts(shape.shapes, out)
